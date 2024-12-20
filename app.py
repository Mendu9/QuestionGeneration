import os
import re
from flask import Flask, request, render_template, flash, redirect, url_for, session
from werkzeug.utils import secure_filename
import PyPDF2
import docx2txt
from pptx import Presentation
import nltk
from nltk.corpus import stopwords
from transformers import T5Tokenizer, T5ForConditionalGeneration, pipeline
import spacy

# Downloads
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)
spacy.cli.download('en_core_web_sm')
nlp = spacy.load('en_core_web_sm')

# Initialize models
summarizer = pipeline('summarization', model='facebook/bart-large-cnn')
tokenizer = T5Tokenizer.from_pretrained("valhalla/t5-base-qg-hl")
qg_model = T5ForConditionalGeneration.from_pretrained("valhalla/t5-base-qg-hl")

app = Flask(__name__)
app.secret_key = 'your_secret_key'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'pptx'}

custom_stopwords = set(stopwords.words('english'))
custom_stopwords.update(['hilke', 'brockmann', 'economies', 'digital', 'societies'])

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def extract_text_from_pdf(file_path):
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        return " ".join([page.extract_text() for page in reader.pages if page.extract_text()])

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return docx2txt.process(file_path)
    elif ext == '.pptx':
        prs = Presentation(file_path)
        return " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

def extract_links(text):
    return re.findall(r'(https?://\S+)', text)

def clean_text(text):
    text = re.sub(r'\n+', ' ', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'http\S+|www\S+', '', text)
    text = re.sub(r'\S+@\S+', '', text)
    text = re.sub(r'[^A-Za-z0-9\s]', '', text)
    return text.strip()

def preprocess_text(text):
    doc = nlp(text)
    words = [token.lemma_ for token in doc if token.is_alpha and not token.is_stop]
    return ' '.join(words)

def summarize_text(text, max_length=130):
    if not text.strip():
        return ""
    chunk_size = 1024
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    summaries = []
    for chunk in chunks:
        s = summarizer(chunk, max_length=max_length, min_length=30, do_sample=False)
        summaries.append(s[0]['summary_text'])
    return ' '.join(summaries).strip()

def chunk_text(text, max_length=300):
    words = text.split()
    if not words:
        return []
    chunks = []
    current_chunk = []
    current_length = 0
    for word in words:
        current_length += len(word) + 1
        if current_length <= max_length:
            current_chunk.append(word)
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_length = len(word)+1
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    return chunks

def generate_questions(text, max_questions=5):
    if not text.strip():
        return []
    input_text = "generate questions: " + text
    input_ids = tokenizer.encode(input_text, return_tensors="pt", max_length=512, truncation=True)
    outputs = qg_model.generate(
        input_ids,
        max_length=128,
        num_return_sequences=min(max_questions, 5),
        num_beams=5,
        early_stopping=True
    )
    questions = [tokenizer.decode(output, skip_special_tokens=True).strip() for output in outputs]
    final_questions = [q for q in questions if len(q.split()) > 3]
    return final_questions

def select_best_questions(questions, count):
    if not questions:
        return []
    unique = list(set(questions))
    unique.sort(key=lambda q: len(q))
    return unique[:count]

@app.route('/', methods=['GET', 'POST'])
def upload_and_process():
    if request.method == 'POST':
        file = request.files.get('file')
        if not file or file.filename == '':
            flash('No file selected')
            return redirect(request.url)
        if allowed_file(file.filename):
            filename = secure_filename(file.filename)
            os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            try:
                original_text = extract_text(file_path)
                text_cleaned = clean_text(original_text)
                processed_text = preprocess_text(text_cleaned)

                summary = summarize_text(processed_text)

                # Generate questions from summary (1 question)
                summary_questions = generate_questions(summary, max_questions=5)
                summary_questions_best = select_best_questions(summary_questions, count=1)

                # Generate questions from chunks (5 questions)
                chunks = chunk_text(processed_text, max_length=300)
                chunk_questions_all = []
                for chunk in chunks:
                    q = generate_questions(chunk, max_questions=5)
                    chunk_questions_all.extend(q)

                summary_set = set(summary_questions_best)
                chunk_filtered = [cq for cq in chunk_questions_all if cq not in summary_set]
                chunk_questions_best = select_best_questions(chunk_filtered, count=5)

                # Combine final questions
                final_questions = summary_questions_best + chunk_questions_best

                session['original_text'] = original_text
                session['summary'] = summary
                session['links'] = extract_links(original_text)
                session['questions'] = final_questions

                return render_template('processed.html',
                                       original_text=original_text,
                                       summary=summary,
                                       links=session['links'])
            except Exception as e:
                flash(f'An error occurred while processing the file: {str(e)}')
                return redirect(request.url)
        else:
            flash('Unsupported file type. Allowed types: pdf, docx, pptx.')
            return redirect(request.url)
    return render_template('upload.html')

@app.route('/questions', methods=['GET'])
def display_questions():
    questions = session.get('questions', [])
    return render_template('questions.html', questions=questions)

@app.route('/grade_single/<int:question_id>', methods=['POST'])
def grade_single(question_id):
    questions = session.get('questions', [])
    original_text = session.get('original_text', '')
    summary = session.get('summary', '')

    if question_id < 0 or question_id >= len(questions):
        flash('Invalid question ID')
        return redirect(url_for('display_questions'))

    user_answer = request.form.get('answer', '').strip()
    question = questions[question_id]

    # Simple grading logic (placeholder)
    context = summary + " " + original_text
    context_words = set(context.lower().split())
    answer_words = set(user_answer.lower().split())

    common = context_words.intersection(answer_words)
    if len(common) > 0:
        feedback = "Good job! Your answer aligns well with the content."
    else:
        feedback = "Try again! Your answer doesn't seem to match the content."

    feedback_key = f'feedback_{question_id}'
    session[feedback_key] = feedback

    return redirect(url_for('display_questions'))

if __name__ == "__main__":
    app.run(debug=True, port=5000)
